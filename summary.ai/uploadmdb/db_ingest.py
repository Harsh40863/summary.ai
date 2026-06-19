# db_ingest.py
import os
import pdfplumber
import fitz  # PyMuPDF
from pymongo import MongoClient
from docx import Document as DocxDocument
from bson import ObjectId
from pymongo.errors import DuplicateKeyError, OperationFailure
from pptx import Presentation
from hashlib import md5
from pymongo import ASCENDING
import io
import re
import unicodedata
from datetime import datetime
from typing import Any, Dict, Optional, Tuple
from dotenv import load_dotenv
from pymongo.collection import Collection

# your embedding function
from embedding.embedder import get_embeddings

def clean_text_for_utf8(text: str) -> str:
    """
    Clean text to ensure it's valid UTF-8 for MongoDB storage.
    Removes or replaces problematic characters.
    """
    if not text:
        return ""
    
    # Normalize unicode characters
    text = unicodedata.normalize('NFKD', text)
    
    # Remove or replace control characters and problematic unicode
    text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)  # Remove control chars
    text = re.sub(r'[\uFFFD]', '', text)  # Remove replacement characters
    text = re.sub(r'[^\x20-\x7E\u00A0-\uFFFF]', ' ', text)  # Replace non-printable with space
    
    # Clean up whitespace
    text = re.sub(r'\s+', ' ', text)  # Multiple spaces to single
    text = text.strip()
    
    return text

# --- MongoDB connection ---
load_dotenv()
uri = os.getenv("MONGODB_URI")
_client = None
_pdf_collection: Optional[Collection[Dict[str, Any]]] = None
_image_collection: Optional[Collection[Dict[str, Any]]] = None


def get_collections() -> Tuple[Collection[Dict[str, Any]], Collection[Dict[str, Any]]]:
    global _client, _pdf_collection, _image_collection
    if _client is None or _pdf_collection is None or _image_collection is None:
        if not uri:
            raise ValueError("MONGODB_URI is not set in .env")
        _client = MongoClient(uri)
        db = _client["hackindia_documents"]
        _pdf_collection = db["pdf_documents"]
        _image_collection = db["images"]

        # Older builds used a global text_hash index, which prevented two users
        # from uploading the same document. Replace it with per-user uniqueness.
        try:
            _pdf_collection.drop_index("text_hash_idx")
        except OperationFailure:
            pass
        try:
            _image_collection.drop_index("image_hash_idx")
        except OperationFailure:
            pass
        _pdf_collection.create_index(
            [("user_email", ASCENDING), ("text_hash", ASCENDING)],
            unique=True,
            name="user_text_hash_idx",
        )
        _image_collection.create_index(
            [("user_email", ASCENDING), ("image_hash", ASCENDING)],
            unique=True,
            name="user_image_hash_idx",
        )

    return _pdf_collection, _image_collection


def _get_uploaded_filename(file_obj: Any) -> str:
    return getattr(file_obj, "filename", None) or getattr(file_obj, "name", "")


def _read_uploaded_file(file_obj: Any) -> bytes:
    if hasattr(file_obj, "file"):
        file_content = file_obj.file.read()
        file_obj.file.seek(0)
        return file_content

    file_content = file_obj.read()
    file_obj.seek(0)
    return file_content


def _document_payload(filename: str, cleaned_text: str, text_hash: str, user_email: Optional[str]) -> Dict[str, Any]:
    return {
        "filename": filename,
        "text": cleaned_text,
        "text_hash": text_hash,
        "embedding": get_embeddings([cleaned_text])[0].tolist(),
        "user_email": user_email,
        "uploaded_at": datetime.utcnow(),
    }


def upload_file_to_db(file_obj: Any, user_email: Optional[str] = None) -> Dict[str, Any]:
    """
    Accepts an uploaded file (FastAPI UploadFile),
    extracts text & images in-memory, and saves them to MongoDB.
    Handles PDF, DOCX, PPTX, TXT, JPG, PNG.
    Returns dict with status and message.
    """
    filename = _get_uploaded_filename(file_obj)
    if not filename:
        return {"text_inserted": False, "images_inserted": 0, "message": "No filename provided"}
    
    ext = filename.lower().split('.')[-1]
    result = {"text_inserted": False, "images_inserted": 0, "message": ""}
    try:
        pdf_collection, image_collection = get_collections()
    except Exception as e:
        result["message"] = f"Database connection failed: {e}"
        return result

    # Read file content into memory
    file_content = _read_uploaded_file(file_obj)

    # ---------- PDFs ----------
    if ext == "pdf":
        # extract text using PyMuPDF (fitz) which is better at preserving spaces
        text = ""
        try:
            doc = fitz.open(stream=file_content, filetype="pdf")
            for page in doc:
                page_text = page.get_text()
                if page_text:
                    text += page_text + "\n"
        except Exception:
            # Fallback
            with pdfplumber.open(io.BytesIO(file_content)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text(x_tolerance=2)
                    if page_text:
                        text += page_text + "\n"

        if text.strip():
            # Clean text for UTF-8 compatibility
            cleaned_text = clean_text_for_utf8(text)
            if not cleaned_text.strip():
                result["message"] += "No valid text content found after cleaning. "
                return result
                
            text_hash = md5(cleaned_text.encode("utf-8")).hexdigest()
            try:
                pdf_collection.insert_one(_document_payload(filename, cleaned_text, text_hash, user_email))
                result["text_inserted"] = True
                result["message"] += "Text saved. "
            except DuplicateKeyError:
                result["message"] += "Duplicate text skipped. "

        # extract images
        doc = fitz.open(stream=file_content, filetype="pdf")
        inserted_images = 0
        for page_index in range(len(doc)):
            page = doc[page_index]
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_hash = md5(image_bytes).hexdigest()
                try:
                    image_collection.insert_one({
                        "filename": filename,
                        "page": page_index + 1,
                        "image_index": img_index,
                        "image_bytes": image_bytes,
                        "image_hash": image_hash,
                        "user_email": user_email,
                        "uploaded_at": datetime.utcnow(),
                    })
                    inserted_images += 1
                except DuplicateKeyError:
                    pass
        result["images_inserted"] = inserted_images
        if inserted_images:
            result["message"] += f"{inserted_images} new image(s) saved."

        return result

    # ---------- DOCX ----------
    elif ext == "docx":
        doc = DocxDocument(io.BytesIO(file_content))
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        if text.strip():
            # Clean text for UTF-8 compatibility
            cleaned_text = clean_text_for_utf8(text)
            if not cleaned_text.strip():
                result["message"] = "No valid text content found after cleaning."
                return result
                
            text_hash = md5(cleaned_text.encode("utf-8")).hexdigest()
            try:
                pdf_collection.insert_one(_document_payload(filename, cleaned_text, text_hash, user_email))
                result["text_inserted"] = True
                result["message"] = "Document saved."
            except DuplicateKeyError:
                result["message"] = "Duplicate document skipped."
        return result

    # ---------- PPTX ----------
    elif ext == "pptx":
        prs = Presentation(io.BytesIO(file_content))
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                shape_text = getattr(shape, "text", "")
                if shape_text.strip():
                    text_runs.append(shape_text)
        text = "\n".join(text_runs)
        if text.strip():
            # Clean text for UTF-8 compatibility
            cleaned_text = clean_text_for_utf8(text)
            if not cleaned_text.strip():
                result["message"] = "No valid text content found after cleaning."
                return result
                
            text_hash = md5(cleaned_text.encode("utf-8")).hexdigest()
            try:
                pdf_collection.insert_one(_document_payload(filename, cleaned_text, text_hash, user_email))
                result["text_inserted"] = True
                result["message"] = "Presentation saved."
            except DuplicateKeyError:
                result["message"] = "Duplicate presentation skipped."
        return result

    # ---------- TXT ----------
    elif ext == "txt":
        text = file_content.decode("utf-8", errors="ignore")
        if text.strip():
            # Clean text for UTF-8 compatibility
            cleaned_text = clean_text_for_utf8(text)
            if not cleaned_text.strip():
                result["message"] = "No valid text content found after cleaning."
                return result
                
            text_hash = md5(cleaned_text.encode("utf-8")).hexdigest()
            try:
                pdf_collection.insert_one(_document_payload(filename, cleaned_text, text_hash, user_email))
                result["text_inserted"] = True
                result["message"] = "Text file saved."
            except DuplicateKeyError:
                result["message"] = "Duplicate text file skipped."
        return result

    # ---------- Images (jpg/png/jpeg) ----------
    elif ext in ["jpg", "jpeg", "png"]:
        image_bytes = file_content
        image_hash = md5(image_bytes).hexdigest()
        try:
            image_collection.insert_one({
                "filename": filename,
                "page": None,
                "image_index": None,
                "image_bytes": image_bytes,
                "image_hash": image_hash,
                "user_email": user_email,
                "uploaded_at": datetime.utcnow(),
            })
            result["images_inserted"] = 1
            result["message"] = "Image saved."
        except DuplicateKeyError:
            result["message"] = "Duplicate image skipped."
        return result

    else:
        result["message"] = "Unsupported file type."
        return result


def fetch_documents_from_db(user_email: Optional[str] = None) -> list[Dict[str, Any]]:
    """
    Fetches all text-based documents from MongoDB and returns them
    in the same structure used by get_documents_from_folder().
    """
    try:
        pdf_collection, _ = get_collections()
    except Exception as e:
        print(f"Warning: Failed to connect to database: {e}")
        return []

    docs = []
    try:
        query = {"user_email": user_email} if user_email else {}
        for i, doc in enumerate(pdf_collection.find(query).sort("uploaded_at", -1)):
            # fallback metadata if not present
            metadata = {
                "author": doc.get("author", "Unknown"),
                "date": doc.get("date", "")  # you can set date at upload time if needed
            }
            docs.append({
                "doc_id": str(doc.get("_id") or doc.get("doc_id", f"doc_{i:03}")),
                "name": doc.get("filename", f"doc_{i:03}"),
                "content": doc.get("text", ""),
                "metadata": metadata,
                "uploaded_at": doc.get("uploaded_at"),
            })
    except Exception as e:
        print(f"Warning: Failed to fetch documents from database: {e}")
        return []
    return docs


def delete_document_from_db(document_id: str, user_email: str) -> bool:
    """Delete one text document owned by user_email."""
    try:
        pdf_collection, image_collection = get_collections()
        doc = pdf_collection.find_one({"_id": ObjectId(document_id), "user_email": user_email})
        if not doc:
            return False
        result = pdf_collection.delete_one({"_id": doc["_id"], "user_email": user_email})
        image_collection.delete_many({"filename": doc.get("filename"), "user_email": user_email})
        return result.deleted_count == 1
    except Exception:
        return False
