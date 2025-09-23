# db_ingest.py
import pdfplumber
import fitz  # PyMuPDF
from pymongo import MongoClient
from docx import Document as DocxDocument
from pymongo.errors import DuplicateKeyError
from pptx import Presentation
from hashlib import md5
from pymongo import ASCENDING
import io
import re
import unicodedata

# your embedding function
from embedding.embedder import get_embeddings

def clean_text_for_utf8(text: str) -> str:
    """
    Clean text to ensure it's valid UTF-8 for MongoDB storage.
    Removes or replaces problematic characters.
    """
    if not text:
        return ""
    
    # Normalize unicode characters
    text = unicodedata.normalize('NFKD', text)
    
    # Remove or replace control characters and problematic unicode
    text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)  # Remove control chars
    text = re.sub(r'[\uFFFD]', '', text)  # Remove replacement characters
    text = re.sub(r'[^\x20-\x7E\u00A0-\uFFFF]', ' ', text)  # Replace non-printable with space
    
    # Clean up whitespace
    text = re.sub(r'\s+', ' ', text)  # Multiple spaces to single
    text = text.strip()
    
    return text

# --- MongoDB connection ---
uri = "mongodb+srv://harshvaishnav0314:Harsh40863@cluster0.aqqqpes.mongodb.net/"
client = MongoClient(uri)
db = client["hackindia_documents"]
pdf_collection = db["pdf_documents"]
image_collection = db["images"]

# Create unique indexes (safe to call repeatedly)
pdf_collection.create_index([("text_hash", ASCENDING)], unique=True, name="text_hash_idx")
image_collection.create_index([("image_hash", ASCENDING)], unique=True, name="image_hash_idx")


def upload_file_to_db(file_obj):
    """
    Accepts an uploaded file (FastAPI UploadFile),
    extracts text & images in-memory, and saves them to MongoDB.
    Handles PDF, DOCX, PPTX, TXT, JPG, PNG.
    Returns dict with status and message.
    """
    # FIXED: Use .filename instead of .name for FastAPI UploadFile
    filename = file_obj.filename
    if not filename:
        return {"text_inserted": False, "images_inserted": 0, "message": "No filename provided"}
    
    ext = filename.lower().split('.')[-1]
    result = {"text_inserted": False, "images_inserted": 0, "message": ""}

    # Read file content into memory
    file_content = file_obj.file.read()
    file_obj.file.seek(0)  # Reset file pointer for potential reuse

    # ---------- PDFs ----------
    if ext == "pdf":
        # extract text
        text = ""
        with pdfplumber.open(io.BytesIO(file_content)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"

        if text.strip():
            # Clean text for UTF-8 compatibility
            cleaned_text = clean_text_for_utf8(text)
            if not cleaned_text.strip():
                result["message"] += "No valid text content found after cleaning. "
                return result
                
            text_hash = md5(cleaned_text.encode("utf-8")).hexdigest()
            embedding = get_embeddings([cleaned_text])[0].tolist()
            try:
                pdf_collection.insert_one({
                    "filename": filename,
                    "text": cleaned_text,
                    "text_hash": text_hash,
                    "embedding": embedding
                })
                result["text_inserted"] = True
                result["message"] += "Text saved. "
            except DuplicateKeyError:
                result["message"] += "Duplicate text skipped. "

        # extract images
        doc = fitz.open(stream=file_content, filetype="pdf")
        inserted_images = 0
        for page_index in range(len(doc)):
            page = doc[page_index]
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_hash = md5(image_bytes).hexdigest()
                try:
                    image_collection.insert_one({
                        "filename": filename,
                        "page": page_index + 1,
                        "image_index": img_index,
                        "image_bytes": image_bytes,
                        "image_hash": image_hash
                    })
                    inserted_images += 1
                except DuplicateKeyError:
                    pass
        result["images_inserted"] = inserted_images
        if inserted_images:
            result["message"] += f"{inserted_images} new image(s) saved."

        return result

    # ---------- DOCX ----------
    elif ext == "docx":
        doc = DocxDocument(io.BytesIO(file_content))
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        if text.strip():
            # Clean text for UTF-8 compatibility
            cleaned_text = clean_text_for_utf8(text)
            if not cleaned_text.strip():
                result["message"] = "No valid text content found after cleaning."
                return result
                
            text_hash = md5(cleaned_text.encode("utf-8")).hexdigest()
            embedding = get_embeddings([cleaned_text])[0].tolist()
            try:
                pdf_collection.insert_one({
                    "filename": filename,
                    "text": cleaned_text,
                    "text_hash": text_hash,
                    "embedding": embedding
                })
                result["text_inserted"] = True
                result["message"] = "Document saved."
            except DuplicateKeyError:
                result["message"] = "Duplicate document skipped."
        return result

    # ---------- PPTX ----------
    elif ext == "pptx":
        prs = Presentation(io.BytesIO(file_content))
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text)
        text = "\n".join(text_runs)
        if text.strip():
            # Clean text for UTF-8 compatibility
            cleaned_text = clean_text_for_utf8(text)
            if not cleaned_text.strip():
                result["message"] = "No valid text content found after cleaning."
                return result
                
            text_hash = md5(cleaned_text.encode("utf-8")).hexdigest()
            embedding = get_embeddings([cleaned_text])[0].tolist()
            try:
                pdf_collection.insert_one({
                    "filename": filename,
                    "text": cleaned_text,
                    "text_hash": text_hash,
                    "embedding": embedding
                })
                result["text_inserted"] = True
                result["message"] = "Presentation saved."
            except DuplicateKeyError:
                result["message"] = "Duplicate presentation skipped."
        return result

    # ---------- TXT ----------
    elif ext == "txt":
        text = file_content.decode("utf-8", errors="ignore")
        if text.strip():
            # Clean text for UTF-8 compatibility
            cleaned_text = clean_text_for_utf8(text)
            if not cleaned_text.strip():
                result["message"] = "No valid text content found after cleaning."
                return result
                
            text_hash = md5(cleaned_text.encode("utf-8")).hexdigest()
            embedding = get_embeddings([cleaned_text])[0].tolist()
            try:
                pdf_collection.insert_one({
                    "filename": filename,
                    "text": cleaned_text,
                    "text_hash": text_hash,
                    "embedding": embedding
                })
                result["text_inserted"] = True
                result["message"] = "Text file saved."
            except DuplicateKeyError:
                result["message"] = "Duplicate text file skipped."
        return result

    # ---------- Images (jpg/png/jpeg) ----------
    elif ext in ["jpg", "jpeg", "png"]:
        image_bytes = file_content
        image_hash = md5(image_bytes).hexdigest()
        try:
            image_collection.insert_one({
                "filename": filename,
                "page": None,
                "image_index": None,
                "image_bytes": image_bytes,
                "image_hash": image_hash
            })
            result["images_inserted"] = 1
            result["message"] = "Image saved."
        except DuplicateKeyError:
            result["message"] = "Duplicate image skipped."
        return result

    else:
        result["message"] = "Unsupported file type."
        return result


def fetch_documents_from_db():
    """
    Fetches all text-based documents from MongoDB and returns them
    in the same structure used by get_documents_from_folder().
    """
    docs = []
    for i, doc in enumerate(pdf_collection.find()):
        # fallback metadata if not present
        metadata = {
            "author": doc.get("author", "Unknown"),
            "date": doc.get("date", "")  # you can set date at upload time if needed
        }
        docs.append({
            "doc_id": doc.get("doc_id", f"doc_{i:03}"),
            "name": doc.get("filename", f"doc_{i:03}"),
            "content": doc.get("text", ""),
            "metadata": metadata
        })
    return docs