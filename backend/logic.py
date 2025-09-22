# ---------------------------
# Core Application Logic
# ---------------------------
import os
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime
from dotenv import load_dotenv
import requests
import nltk
import numpy as np
from nltk.tokenize import sent_tokenize
from sklearn.cluster import KMeans
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from sklearn.metrics.pairwise import cosine_similarity
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO

# Custom modules
from text_preprocessing.cleaner import clean_text
from embedding.embedder import get_embeddings
from search.vector_search import VectorSearchEngine
from summarization.summarizer import flan_summarize_with_query as summarize_text_with_query
from rag.rag_engine import generate_answer_with_rag
from rag.langchain_tools import GoogleSearchTool
from uploadmdb.db_ingest import upload_file_to_db, fetch_documents_from_db

# ---------------------------
# Load environment variables
# ---------------------------
load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
GOOGLE_API_KEY2 = os.getenv("GOOGLE_API_KEY2")
GOOGLE_CSE_ID = os.getenv("GOOGLE_CSE_ID")
UNSPLASH_ACCESS_KEY = os.getenv("UNSPLASH_ACCESS_KEY")

class DocumentSearchEngine:
    """
    Core logic class for document search and analysis functionality
    """
    
    def __init__(self):
        self.documents = []
        self.doc_texts = []
        self.doc_embeddings = None
        self.index_to_docinfo = {}
        self.topic_labels = None
        self.search_engine = None
        self.n_clusters = 3
        self._initialize_nltk()
        self._load_documents()
        
    def _initialize_nltk(self):
        """Download required NLTK data"""
        try:
            nltk.data.find("tokenizers/punkt")
        except LookupError:
            nltk.download("punkt")
    
    def _load_documents(self):
        """Load and preprocess documents from database"""
        self.documents = fetch_documents_from_db()
        
        if not self.documents:
            raise ValueError("No valid documents found in the database.")
        
        # Process documents
        self.doc_texts = [clean_text(doc["content"]) for doc in self.documents]
        self.doc_embeddings = get_embeddings(self.doc_texts)
        self.index_to_docinfo = {i: doc for i, doc in enumerate(self.documents)}
        
        # Create clusters
        self.n_clusters = min(len(self.documents), 3)
        kmeans = KMeans(n_clusters=self.n_clusters, random_state=42).fit(self.doc_embeddings)
        self.topic_labels = kmeans.labels_
        
        # Initialize search engine
        self.search_engine = VectorSearchEngine(dimension=self.doc_embeddings.shape[1])
        self.search_engine.build_index(self.doc_embeddings, self.doc_texts)
    
    def reload_documents(self):
        """Reload documents from database (useful after new uploads)"""
        self._load_documents()
    
    def upload_files(self, uploaded_files: List[Any]) -> List[Dict[str, str]]:
        """
        Upload multiple files to database
        
        Args:
            uploaded_files: List of file objects
            
        Returns:
            List of upload results with filename and status
        """
        results = []
        for uploaded_file in uploaded_files:
            try:
                result = upload_file_to_db(uploaded_file)
                results.append({
                    "filename": uploaded_file.filename if hasattr(uploaded_file, 'filename') else str(uploaded_file),
                    "status": "success",
                    "message": result
                })
            except Exception as e:
                results.append({
                    "filename": uploaded_file.filename if hasattr(uploaded_file, 'filename') else str(uploaded_file),
                    "status": "error",
                    "message": str(e)
                })
        
        # Reload documents after upload
        self.reload_documents()
        return results
    
    def google_search(self, query: str, api_key: str = None, cse_id: str = None, num: int = 5) -> Dict[str, Any]:
        """
        Perform Google Custom Search
        
        Args:
            query: Search query string
            api_key: Google API key (uses env variable if not provided)
            cse_id: Custom Search Engine ID (uses env variable if not provided)
            num: Number of results to return
            
        Returns:
            Dictionary with search results or error
        """
        api_key = api_key or GOOGLE_API_KEY
        cse_id = cse_id or GOOGLE_CSE_ID
        
        url = "https://www.googleapis.com/customsearch/v1"
        params = {"q": query, "key": api_key, "cx": cse_id, "num": num}
        
        try:
            response = requests.get(url, params=params)
            response.raise_for_status()
            return response.json()
        except requests.exceptions.RequestException as e:
            return {"error": str(e)}
    
    def find_relevant_documents(self, query: str, threshold: float = 0.80) -> List[Dict[str, Any]]:
        """
        Find documents relevant to the query
        
        Args:
            query: Search query string
            threshold: Similarity threshold for relevance
            
        Returns:
            List of relevant document results sorted by relevance score
        """
        clean_query = clean_text(query)
        query_embedding = get_embeddings([clean_query])
        
        relevant_results = []
        
        for idx, doc in enumerate(self.documents):
            sentences = sent_tokenize(doc["content"])
            if not sentences:
                continue
            
            sentence_embeddings = get_embeddings(sentences)
            # Use cosine similarity instead of dot product
            sim_scores = cosine_similarity(sentence_embeddings, query_embedding).flatten()
            best_score = np.max(sim_scores)
            best_idx = np.argmax(sim_scores)
            
            if best_score < threshold:
                continue
            
            # Get relevant context chunk
            start = max(0, best_idx - 1)
            end = min(len(sentences), best_idx + 2)
            relevant_chunk = " ".join(sentences[start:end])
            
            # Generate query-specific summary
            query_summary = generate_answer_with_rag(query, relevant_chunk, doc["name"])
            
            relevant_results.append({
                "doc": doc,
                "best_sentence": sentences[best_idx],
                "score": float(best_score),  # Convert to Python float
                "relevant_chunk": relevant_chunk,
                "summary": query_summary,
                "doc_index": idx
            })
        
        # Sort by score (highest first)
        # Return only top 5 most relevant results
        sorted_results = sorted(relevant_results, key=lambda x: x["score"], reverse=True)
        return sorted_results[:5]
    
    def search_action(self, query: str) -> Dict[str, Any]:
        """
        Perform search action - returns document summaries
        
        Args:
            query: Search query string
            
        Returns:
            Dictionary with search results and metadata
        """
        relevant_results = self.find_relevant_documents(query)
        
        if not relevant_results:
            return {
                "success": False,
                "message": "No relevant documents found for your query.",
                "results": []
            }
        
        search_results = []
        for result in relevant_results:
            doc = result["doc"]
            search_results.append({
                "document_name": doc["name"],
                "best_sentence": result["best_sentence"],
                "score": result["score"],
                "summary": result["summary"],
                "action_type": "search"
            })
        
        return {
            "success": True,
            "message": f"Found {len(search_results)} relevant documents",
            "results": search_results,
            "action": "search"
        }
    
    def explore_action(self, query: str) -> Dict[str, Any]:
        """
        Perform explore action - combines document insights with web search
        
        Args:
            query: Search query string
            
        Returns:
            Dictionary with exploration results
        """
        relevant_results = self.find_relevant_documents(query)
        
        if not relevant_results:
            return {
                "success": False,
                "message": "No relevant documents found for your query.",
                "results": []
            }
        
        explore_results = []
        for result in relevant_results:
            doc = result["doc"]
            summary = result["summary"]
            
            # Enhanced query for web search
            enhanced_query = f"{query} - Context: {summary[:250]}"
            
            try:
                search_tool = GoogleSearchTool(api_key=GOOGLE_API_KEY2, cse_id=GOOGLE_CSE_ID)
                web_content = search_tool.run(enhanced_query)
                
                explore_results.append({
                    "document_name": doc["name"],
                    "best_sentence": result["best_sentence"],
                    "score": result["score"],
                    "web_content": web_content,
                    "action_type": "explore"
                })
            except Exception as e:
                explore_results.append({
                    "document_name": doc["name"],
                    "best_sentence": result["best_sentence"],
                    "score": result["score"],
                    "web_content": f"Error fetching web content: {str(e)}",
                    "action_type": "explore",
                    "error": True
                })
        
        return {
            "success": True,
            "message": f"Explored {len(explore_results)} relevant documents",
            "results": explore_results,
            "action": "explore"
        }
    
    def think_action(self, query: str) -> Dict[str, Any]:
        """
        Perform think action - generates refined insights using Google Gen AI
        
        Args:
            query: Search query string
            
        Returns:
            Dictionary with thinking results
        """
        relevant_results = self.find_relevant_documents(query)
        
        if not relevant_results:
            return {
                "success": False,
                "message": "No relevant documents found for your query.",
                "results": []
            }
        
        think_results = []
        for result in relevant_results:
            doc = result["doc"]
            summary = result["summary"]
            
            try:
                from rag.Google_Gen_AI import generate_answer_with_google
                refined_answer = generate_answer_with_google(query, summary, doc["name"])
                
                think_results.append({
                    "document_name": doc["name"],
                    "best_sentence": result["best_sentence"],
                    "score": result["score"],
                    "refined_insight": refined_answer,
                    "action_type": "think"
                })
            except Exception as e:
                think_results.append({
                    "document_name": doc["name"],
                    "best_sentence": result["best_sentence"],
                    "score": result["score"],
                    "refined_insight": f"Error during insight generation: {str(e)}",
                    "action_type": "think",
                    "error": True
                })
        
        return {
            "success": True,
            "message": f"Generated insights for {len(think_results)} relevant documents",
            "results": think_results,
            "action": "think"
        }
    
    def process_query(self, query: str, action: str) -> Dict[str, Any]:
        """
        Main method to process query with specified action
        
        Args:
            query: Search query string
            action: Action type ("search", "explore", "think", "ppt")
            
        Returns:
            Dictionary with results based on action type
        """
        if not query or not query.strip():
            return {
                "success": False,
                "message": "Query cannot be empty",
                "results": []
            }
        
        action = action.lower().strip()
        
        if action == "search":
            return self.search_action(query)
        elif action == "explore":
            return self.explore_action(query)
        elif action == "think":
            return self.think_action(query)
        elif action == "ppt":
            return self.generate_ppt_action(query)
        else:
            return {
                "success": False,
                "message": f"Unknown action: {action}. Available actions: search, explore, think, ppt",
                "results": []
            }

    def _slugify_filename(self, text: str) -> str:
        """Create a filesystem-safe slug for filenames"""
        import re
        slug = re.sub(r"[^a-zA-Z0-9\-_. ]+", "", text).strip().replace(" ", "_")
        if not slug:
            slug = "presentation"
        return slug[:60]

    def _create_decorative_presentation(self, title: str, summary: str) -> str:
        """
        Create a decorative PPTX file from title and summary.

        Returns the absolute file path to the generated PPTX.
        """
        # Create presentation with a widescreen layout
        prs = Presentation()

        # Simple brand palette
        primary = RGBColor(30, 144, 255)      # DodgerBlue
        secondary = RGBColor(230, 245, 255)   # Light Blue BG
        accent = RGBColor(76, 175, 80)        # Green

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        subtitle = slide.placeholders[1]
        subtitle.text = f"Auto-generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        # Style title
        title_tf = slide.shapes.title.text_frame
        title_tf.paragraphs[0].font.size = Pt(44)
        title_tf.paragraphs[0].font.bold = True

        # Summary bullets slide
        bullet_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_layout)
        slide2.shapes.title.text = "Executive Summary"
        body = slide2.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear()
        # Split summary into bullet points by sentences/newlines
        points = [p.strip() for p in summary.replace("\n", " ").split('.') if p.strip()]
        if not points:
            points = [summary]
        for i, point in enumerate(points[:8]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.level = 0
        # Add a decorative band under the title
        band_shape = slide2.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.7), Inches(9.0), Inches(0.15)
        )
        band_fill = band_shape.fill
        band_fill.solid()
        band_fill.fore_color.rgb = primary

        # Decorative content slides for first few points
        # Fetch 1 hero image once and reuse (speed optimization)
        hero_image = self._fetch_image_bytes(f"{title} concept illustration")

        for idx, point in enumerate(points[:5], start=1):
            content_layout = prs.slide_layouts[5]  # Title Only
            s = prs.slides.add_slide(content_layout)
            s.shapes.title.text = f"Key Point {idx}"
            # Add a rounded rectangle with the point text
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(2)
            shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = secondary
            line = shape.line
            line.color.rgb = primary
            text_frame = shape.text_frame
            text_frame.text = point
            for para in text_frame.paragraphs:
                para.font.size = Pt(20)

            # Add a small decorative circle as an icon
            circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(1.8), Inches(0.6), Inches(0.6))
            circle.fill.solid()
            circle.fill.fore_color.rgb = primary
            circle.line.color.rgb = primary

            # Try to add a relevant image on the right
            image_bytes = hero_image  # reuse hero image for speed
            if image_bytes is not None:
                try:
                    pic_left = Inches(6.8)
                    pic_top = Inches(1.8)
                    pic_width = Inches(2.2)
                    s.shapes.add_picture(BytesIO(image_bytes), pic_left, pic_top, width=pic_width)
                except Exception:
                    # Ignore image placement errors
                    pass

        # Ensure output directory exists
        out_dir = os.path.join(os.getcwd(), "generated")
        os.makedirs(out_dir, exist_ok=True)
        file_name = f"{self._slugify_filename(title)}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        file_path = os.path.join(out_dir, file_name)
        prs.save(file_path)
        return file_path

    def _fetch_image_bytes(self, query: str) -> Any:
        """
        Fetch a relevant image for the query from Unsplash. Returns bytes or None.
        """
        try:
            if not UNSPLASH_ACCESS_KEY:
                return None
            url = "https://api.unsplash.com/search/photos"
            params = {
                "query": query,
                "per_page": 1,
                "orientation": "landscape"
            }
            headers = {"Authorization": f"Client-ID {UNSPLASH_ACCESS_KEY}"}
            resp = requests.get(url, params=params, headers=headers, timeout=3)
            resp.raise_for_status()
            data = resp.json()
            if not data.get("results"):
                return None
            img_url = data["results"][0]["urls"].get("regular") or data["results"][0]["urls"].get("small")
            if not img_url:
                return None
            img_resp = requests.get(img_url, timeout=3)
            img_resp.raise_for_status()
            return img_resp.content
        except Exception:
            return None

    def _generate_ppt_outline(self, query: str, summary: str) -> Dict[str, Any]:
        """
        Use Google Gen AI to turn the summary into a clean outline for PPT creation.
        Returns a dict with structure: { title, agenda[], sections:[{heading, bullets[], image_keywords[]}], conclusion[] }
        Falls back to a heuristic outline if the AI call fails.
        """
        # Default heuristic
        def heuristic_outline() -> Dict[str, Any]:
            sentences = [s.strip() for s in summary.replace("\n", " ").split('.') if s.strip()] or [summary]
            top_points = sentences[:6]
            sections = []
            for idx, chunk in enumerate([top_points[:2], top_points[2:4], top_points[4:6]]):
                if not chunk:
                    continue
                sections.append({
                    "heading": f"Section {idx+1}",
                    "bullets": chunk,
                    "image_keywords": [query]
                })
            return {
                "title": query,
                "agenda": [s[:90] for s in top_points[:4]],
                "sections": sections,
                "conclusion": ["Key takeaways:"] + [s[:100] for s in top_points[:2]]
            }

        try:
            from rag.Google_Gen_AI import generate_answer_with_google
            prompt = (
                "You are a presentation assistant. Given a topic and context summary,"
                " produce a concise JSON for slides with keys: title (string), agenda (array of 3-5 short items),"
                " sections (array of 2-4 objects each with: heading (string), bullets (array of 3-5 short items),"
                " image_keywords (array of 1-3 short keywords)), conclusion (array of 2-4 short items)."
                " Keep text concise, no markdown or extra commentary."
                f"\nTopic: {query}\nContext: {summary}"
            )
            raw = generate_answer_with_google("Create PPT outline JSON", prompt, "ppt-outline")
            import json
            # Try to extract JSON from raw
            start = raw.find('{')
            end = raw.rfind('}')
            obj = json.loads(raw[start:end+1]) if start != -1 and end != -1 else json.loads(raw)
            # Validate minimal structure
            if not isinstance(obj, dict) or "title" not in obj:
                return heuristic_outline()
            obj.setdefault("agenda", [])
            obj.setdefault("sections", [])
            obj.setdefault("conclusion", [])
            return obj
        except Exception:
            return heuristic_outline()

    def _create_presentation_from_outline(self, outline: Dict[str, Any]) -> str:
        """Create a polished PPTX using a structured outline."""
        prs = Presentation()

        primary = RGBColor(30, 144, 255)
        secondary = RGBColor(230, 245, 255)
        dark = RGBColor(33, 33, 33)

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = outline.get("title") or "Presentation"
        # Improve title visibility
        try:
            title_tf = slide.shapes.title.text_frame
            if title_tf.paragraphs:
                title_tf.paragraphs[0].font.bold = True
                title_tf.paragraphs[0].font.size = Pt(44)
                title_tf.paragraphs[0].font.color.rgb = dark
        except Exception:
            pass
        subtitle = slide.placeholders[1]
        subtitle.text = "Auto-generated deck"
        try:
            sub_tf = subtitle.text_frame
            if sub_tf.paragraphs:
                sub_tf.paragraphs[0].font.color.rgb = dark
                sub_tf.paragraphs[0].font.size = Pt(18)
        except Exception:
            pass
        # Decorative band
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5), Inches(10), Inches(0.4))
        band.fill.solid()
        band.fill.fore_color.rgb = primary

        # Agenda slide
        if outline.get("agenda"):
            layout = prs.slide_layouts[1]
            s2 = prs.slides.add_slide(layout)
            s2.shapes.title.text = "Agenda"
            try:
                if s2.shapes.title.text_frame.paragraphs:
                    s2.shapes.title.text_frame.paragraphs[0].font.color.rgb = dark
                    s2.shapes.title.text_frame.paragraphs[0].font.bold = True
            except Exception:
                pass
            body = s2.shapes.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, item in enumerate(outline["agenda"][:6]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = str(item)
                p.level = 0
                try:
                    p.font.color.rgb = dark
                    p.font.size = Pt(18)
                except Exception:
                    pass
            # band
            band2 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.7), Inches(9.0), Inches(0.15))
            band2.fill.solid()
            band2.fill.fore_color.rgb = primary

        # Prepare hero image once for speed
        hero_image = None
        # Try to derive a topic keyword
        topic_kw = outline.get("title") or ""
        hero_image = self._fetch_image_bytes(f"{topic_kw} concept illustration")

        # Section slides
        for section in outline.get("sections", [])[:4]:
            layout = prs.slide_layouts[5]  # Title Only
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = section.get("heading") or "Section"
            try:
                if s.shapes.title.text_frame.paragraphs:
                    s.shapes.title.text_frame.paragraphs[0].font.color.rgb = dark
                    s.shapes.title.text_frame.paragraphs[0].font.bold = True
            except Exception:
                pass

            # Text container
            left = Inches(0.8)
            top = Inches(1.6)
            width = Inches(5.8)
            height = Inches(3.6)
            rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            rect.fill.solid()
            rect.fill.fore_color.rgb = secondary
            rect.line.color.rgb = primary
            tf = rect.text_frame
            tf.clear()
            bullets = section.get("bullets", [])
            for i, b in enumerate(bullets[:6]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = str(b)
                p.level = 0
                try:
                    p.font.color.rgb = dark
                    p.font.size = Pt(18)
                except Exception:
                    pass

            # Images: try 1-2 images; reuse hero if available; otherwise fetch using keywords
            images_added = 0
            img_left = Inches(6.8)
            img_top = Inches(1.4)
            img_width = Inches(2.2)

            if hero_image is not None:
                try:
                    s.shapes.add_picture(BytesIO(hero_image), img_left, img_top, width=img_width)
                    images_added += 1
                except Exception:
                    pass

            if images_added < 2:
                kw_list = section.get("image_keywords") or [topic_kw]
                for kw in kw_list:
                    img = self._fetch_image_bytes(str(kw))
                    if img is not None:
                        try:
                            s.shapes.add_picture(BytesIO(img), img_left, img_top + Inches(2.4 * images_added), width=img_width)
                            images_added += 1
                            if images_added >= 2:
                                break
                        except Exception:
                            continue

            # Decorative icon
            icon = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.4), Inches(1.2), Inches(0.5), Inches(0.5))
            icon.fill.solid()
            icon.fill.fore_color.rgb = primary
            icon.line.color.rgb = primary

        # Conclusion slide
        if outline.get("conclusion"):
            layout = prs.slide_layouts[1]
            sc = prs.slides.add_slide(layout)
            sc.shapes.title.text = "Conclusion"
            try:
                if sc.shapes.title.text_frame.paragraphs:
                    sc.shapes.title.text_frame.paragraphs[0].font.color.rgb = dark
                    sc.shapes.title.text_frame.paragraphs[0].font.bold = True
            except Exception:
                pass
            body = sc.shapes.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, item in enumerate(outline["conclusion"][:6]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = str(item)
                p.level = 0
                try:
                    p.font.color.rgb = dark
                    p.font.size = Pt(18)
                except Exception:
                    pass
            band3 = sc.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.7), Inches(9.0), Inches(0.15))
            band3.fill.solid()
            band3.fill.fore_color.rgb = primary

        out_dir = os.path.join(os.getcwd(), "generated")
        os.makedirs(out_dir, exist_ok=True)
        file_name = f"{self._slugify_filename(outline.get('title') or 'presentation')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        file_path = os.path.join(out_dir, file_name)
        prs.save(file_path)
        return file_path

    def generate_ppt_action(self, query: str) -> Dict[str, Any]:
        """
        Generate a decorative PPT using the best RAG summary for the query.

        Returns a dict containing path and metadata of the generated PPT.
        """
        relevant_results = self.find_relevant_documents(query)

        if not relevant_results:
            return {
                "success": False,
                "message": "No relevant documents found for your query.",
                "results": []
            }

        # Use the top result's summary to build the deck
        top = relevant_results[0]
        doc = top["doc"]
        summary = top["summary"]

        try:
            # Generate structured outline first for better quality content
            outline = self._generate_ppt_outline(query=query, summary=summary)
            ppt_path = self._create_presentation_from_outline(outline)
            return {
                "success": True,
                "message": f"PPT generated successfully from document: {doc['name']}",
                "results": [{
                    "document_name": doc["name"],
                    "score": top["score"],
                    "ppt_path": ppt_path,
                    "file_name": os.path.basename(ppt_path),
                    "action_type": "ppt"
                }],
                "action": "ppt"
            }
        except Exception as e:
            return {
                "success": False,
                "message": f"Failed to generate PPT: {str(e)}",
                "results": []
            }
    
    def get_documents_info(self) -> Dict[str, Any]:
        """
        Get information about loaded documents
        
        Returns:
            Dictionary with document statistics and info
        """
        if not self.documents:
            return {
                "total_documents": 0,
                "documents": [],
                "clusters": 0
            }
        
        doc_info = []
        for doc in self.documents:
            doc_info.append({
                "name": doc["name"],
                "content_length": len(doc["content"]),
                "upload_date": doc.get("upload_date", "Unknown")
            })
        
        return {
            "total_documents": len(self.documents),
            "documents": doc_info,
            "clusters": self.n_clusters
        }
    
    def health_check(self) -> Dict[str, Any]:
        """
        Check if the system is ready to process queries
        
        Returns:
            Dictionary with system health status
        """
        try:
            has_docs = len(self.documents) > 0
            has_embeddings = self.doc_embeddings is not None
            has_search_engine = self.search_engine is not None
            
            return {
                "status": "healthy" if all([has_docs, has_embeddings, has_search_engine]) else "unhealthy",
                "has_documents": has_docs,
                "has_embeddings": has_embeddings,
                "has_search_engine": has_search_engine,
                "document_count": len(self.documents),
                "timestamp": datetime.now().isoformat()
            }
        except Exception as e:
            return {
                "status": "error",
                "error": str(e),
                "timestamp": datetime.now().isoformat()
            }


# Global instance for reuse
_search_engine_instance = None

def get_search_engine() -> DocumentSearchEngine:
    """
    Get or create global search engine instance
    
    Returns:
        DocumentSearchEngine instance
    """
    global _search_engine_instance
    if _search_engine_instance is None:
        _search_engine_instance = DocumentSearchEngine()
    return _search_engine_instance

def reset_search_engine():
    """Reset global search engine instance (useful for testing or reinitialization)"""
    global _search_engine_instance
    _search_engine_instance = None